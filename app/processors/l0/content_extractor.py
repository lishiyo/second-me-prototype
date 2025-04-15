import io
from typing import Optional, Dict, Any

import docx
import PyPDF2
from pptx import Presentation

from app.processors.l0.models import FileInfo
from app.processors.l0.utils import retry, safe_execute, setup_logger

# Set up logger
logger = setup_logger(__name__)

class ContentExtractor:
    """
    Extracts text content from various file types.
    Supports: PDF, DOCX, PPTX, TXT and other text-based formats.
    """
    
    def __init__(self):
        self.extractors = {
            'application/pdf': self._extract_pdf,
            'application/vnd.openxmlformats-officedocument.wordprocessingml.document': self._extract_docx,
            'application/vnd.openxmlformats-officedocument.presentationml.presentation': self._extract_pptx,
            'text/plain': self._extract_text,
            'text/markdown': self._extract_text,
            'text/html': self._extract_text,
            'application/json': self._extract_text,
        }
    
    def extract(self, file_info: FileInfo) -> str:
        """
        Extract content from a file based on its content type.
        
        Args:
            file_info: FileInfo object containing file metadata and content
            
        Returns:
            Extracted text content
            
        Raises:
            ValueError: If file content is None or content type is unsupported
        """
        if file_info.content is None:
            raise ValueError("File content is None")
        
        extractor = self.extractors.get(file_info.content_type)
        
        if extractor is None:
            logger.warning(f"No specific extractor for content type: {file_info.content_type}. Trying text extractor.")
            extractor = self._extract_text
        
        logger.info(f"Extracting content from {file_info.filename} ({file_info.content_type})")
        return extractor(file_info)
    
    @retry(max_retries=2)
    def _extract_pdf(self, file_info: FileInfo) -> str:
        """Extract text from PDF files."""
        pdf_content = []
        
        with io.BytesIO(file_info.content) as pdf_file:
            reader = PyPDF2.PdfReader(pdf_file)
            
            for page_num in range(len(reader.pages)):
                page = reader.pages[page_num]
                text = safe_execute(page.extract_text, default_value="")
                
                if text:
                    pdf_content.append(text)
                else:
                    logger.warning(f"Could not extract text from page {page_num} in {file_info.filename}")
        
        return "\n\n".join(pdf_content)
    
    def _extract_docx(self, file_info: FileInfo) -> str:
        """Extract text from DOCX files."""
        with io.BytesIO(file_info.content) as docx_file:
            doc = docx.Document(docx_file)
            return "\n".join([paragraph.text for paragraph in doc.paragraphs if paragraph.text])
    
    def _extract_pptx(self, file_info: FileInfo) -> str:
        """Extract text from PPTX files."""
        content = []
        
        with io.BytesIO(file_info.content) as pptx_file:
            presentation = Presentation(pptx_file)
            
            for i, slide in enumerate(presentation.slides, 1):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text)
                
                if slide_text:
                    content.append(f"Slide {i}:\n" + "\n".join(slide_text))
        
        return "\n\n".join(content)
    
    def _extract_text(self, file_info: FileInfo) -> str:
        """Extract text from plain text files."""
        content = file_info.content
        
        # Handle different content types
        if isinstance(content, bytes):
            try:
                return content.decode('utf-8')
            except UnicodeDecodeError:
                logger.warning(f"UTF-8 decoding failed for {file_info.filename}, trying with errors='replace'")
                return content.decode('utf-8', errors='replace')
        elif isinstance(content, str):
            return content
        else:
            # For any other type, convert to string
            logger.warning(f"Unexpected content type: {type(content)}, converting to string")
            return str(content)

    def extract_content(self, content: Any, content_type: str) -> str:
        """
        Extract content from binary data based on content type.
        This method is used by the document processor for direct binary content.
        
        Args:
            content: Binary file content or string content
            content_type: MIME type of the content
            
        Returns:
            Extracted text content as a string
        """
        if content is None:
            raise ValueError("Content is None")
            
        # If it's already a string, return it directly
        if isinstance(content, str):
            return content
            
        # Create a temporary FileInfo object to reuse our extraction methods
        temp_file_info = FileInfo(
            document_id="temp",
            filename="temp",
            content_type=content_type,
            s3_path="",
            content=content
        )
        
        # Use the appropriate extractor based on content type
        extractor = self.extractors.get(content_type)
        
        if extractor is None:
            logger.warning(f"No specific extractor for content type: {content_type}. Trying text extractor.")
            extractor = self._extract_text
            
        logger.info(f"Extracting content from binary data ({content_type})")
        
        try:
            return extractor(temp_file_info)
        except Exception as e:
            logger.error(f"Error extracting content: {str(e)}")
            # If extraction fails, try a basic UTF-8 decode for binary data
            if isinstance(content, bytes):
                try:
                    return content.decode('utf-8', errors='replace')
                except Exception:
                    return f"[Error extracting content: {str(e)}]"
            else:
                return f"[Error extracting content: {str(e)}]"
